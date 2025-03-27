import requests
import matplotlib.pyplot as plt
from datetime import datetime
import os
from pptx import Presentation
from pptx.util import Inches

def get_weather_data():
    api_key = "c4a1259f23msh4e1b8e6b6d4f2e7p1f6b3bjsn2f4a3c0b3f2d"
    url = "https://weatherapi-com.p.rapidapi.com/forecast.json"
    
    querystring = {
        "q": "Moscow",
        "days": "3"
    }
    
    headers = {
        "X-RapidAPI-Key": api_key,
        "X-RapidAPI-Host": "weatherapi-com.p.rapidapi.com"
    }
    
    try:
        response = requests.get(url, headers=headers, params=querystring)
        response.raise_for_status()
        return response.json()
    except Exception as e:
        print(f"Error fetching weather data: {str(e)}")
        return None

def create_current_weather_plot(data, filename):
    if not data or 'current' not in data:
        return
    
    plt.figure(figsize=(12, 6))
    
    current = data['current']
    temp_c = current['temp_c']
    feelslike_c = current['feelslike_c']
    
    plt.bar(['Температура', 'Ощущается как'], [temp_c, feelslike_c])
    plt.title(f"Текущая погода в Москве\n{datetime.now().strftime('%d.%m.%Y %H:%M')}")
    plt.ylabel('Температура (°C)')
    
    plt.grid(True)
    plt.tight_layout()
    plt.savefig(filename)
    plt.close()

def create_hourly_forecast_plot(data, period_name, hours, filename):
    if not data or 'forecast' not in data:
        return
    
    plt.figure(figsize=(12, 6))
    
    times = []
    temps = []
    
    for hour in data['forecast']['forecastday'][0]['hour'][hours[0]:hours[1]]:
        time = datetime.fromisoformat(hour['time']).strftime('%H:%M')
        times.append(time)
        temps.append(hour['temp_c'])
    
    plt.plot(times, temps, marker='o')
    plt.title(f"Прогноз на {period_name}")
    plt.xlabel('Время')
    plt.ylabel('Температура (°C)')
    
    plt.grid(True)
    plt.xticks(rotation=45)
    plt.tight_layout()
    
    plt.savefig(filename)
    plt.close()

def create_presentation(images):
    prs = Presentation()
    
    # Add a title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    title.text = "Погода в Москве"
    
    # Add content slide with weather plots
    content_slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    # Position the plots in a 2x2 grid
    positions = [
        (Inches(0.5), Inches(0.5), Inches(4.5), Inches(3)),
        (Inches(5), Inches(0.5), Inches(4.5), Inches(3)),
        (Inches(0.5), Inches(3.5), Inches(4.5), Inches(3)),
        (Inches(5), Inches(3.5), Inches(4.5), Inches(3))
    ]
    
    for image, position in zip(images, positions):
        if os.path.exists(image):
            content_slide.shapes.add_picture(
                image,
                position[0], position[1],
                width=position[2], height=position[3]
            )
    
    prs.save('weather_report.pptx')

def main():
    if not os.path.exists('weather_plots'):
        os.makedirs('weather_plots')
    
    # Get weather data
    weather_data = get_weather_data()
    
    if weather_data:
        plots = []
        
        # Current weather
        current_plot = 'weather_plots/weather_current.png'
        create_current_weather_plot(weather_data, current_plot)
        plots.append(current_plot)
        
        # Forecast for different parts of the day
        periods = [
            ('утро', (6, 12)),
            ('день', (12, 18)),
            ('вечер', (18, 24))
        ]
        
        for period_name, hours in periods:
            filename = f'weather_plots/weather_forecast_{period_name}.png'
            create_hourly_forecast_plot(weather_data, period_name, hours, filename)
            plots.append(filename)
        
        # Create presentation
        create_presentation(plots)
        print("Презентация успешно создана!")
    else:
        print("Не удалось получить данные о погоде")

if __name__ == "__main__":
    main() 