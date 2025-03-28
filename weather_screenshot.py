from playwright.sync_api import sync_playwright
from PIL import Image
import time
import random
from pptx import Presentation
from pptx.util import Inches
import os
from playwright.async_api import async_playwright

def simulate_human_behavior(page):
    # Случайная задержка
    time.sleep(random.uniform(2, 4))
    
    # Случайное движение мыши
    page.mouse.move(random.randint(100, 500), random.randint(100, 500))
    
    # Небольшой скролл
    page.mouse.wheel(0, random.randint(100, 300))
    time.sleep(random.uniform(1, 2))

async def take_screenshot(page, url, output_path):
    try:
        # Faster navigation with shorter timeout
        await page.goto(url, wait_until='domcontentloaded', timeout=20000)
        
        # Wait for the main content to load
        content_selector = '.content-column[data-column="C1"]'
        await page.wait_for_selector(content_selector, state='visible', timeout=20000)
        
        # Remove ads and unnecessary elements
        await page.evaluate("""() => {
            const ads = document.querySelectorAll('.ad-content, .ad-block, .banner');
            ads.forEach(ad => ad.remove());
        }""")
        
        # Get the content element and precipitation element bounding boxes
        content = page.locator(content_selector)
        precipitation = page.get_by_text("Осадки в жидком эквиваленте, мм")
        
        content_box = await content.bounding_box()
        precip_box = await precipitation.bounding_box()
        
        if content_box and precip_box:
            # Calculate the crop box
            box = {
                'x': max(0, content_box['x'] - 20),
                'y': max(0, content_box['y'] - 15),
                'width': content_box['width'] + 40,
                'height': (precip_box['y'] + precip_box['height'] + 75 - content_box['y'])
            }
            
            # Take screenshot directly with the calculated box
            await page.screenshot(
                path=output_path,
                clip=box
            )
            
            print(f"Screenshot saved to {output_path}")
            return True
        else:
            print(f"Could not get bounding box for content on {url}")
            return False
    except Exception as e:
        print(f"Error taking screenshot of {url}: {str(e)}")
        return False

async def create_presentation(screenshots):
    prs = Presentation()
    
    # Создаем слайд
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 6 - пустой макет
    
    # Размеры слайда в дюймах (стандартный размер 10x7.5)
    slide_width = Inches(10)
    slide_height = Inches(7.5)
    
    # Размеры для каждого скриншота (чуть меньше половины слайда)
    width = Inches(4.8)  # Увеличили ширину
    height = Inches(3.3)  # Увеличили высоту
    
    # Отступы
    margin = Inches(0.2)  # Маленький отступ между скриншотами
    
    # Позиции для сетки 2x2 (автоматический расчет для центрирования)
    left_x = (slide_width - (2 * width + margin)) / 2  # Центрируем по горизонтали
    top_y = (slide_height - (2 * height + margin)) / 2  # Центрируем по вертикали
    
    positions = [
        # Верхний ряд
        {'left': left_x, 'top': top_y},  # Левый верхний
        {'left': left_x + width + margin, 'top': top_y},  # Правый верхний
        # Нижний ряд
        {'left': left_x, 'top': top_y + height + margin},  # Левый нижний
        {'left': left_x + width + margin, 'top': top_y + height + margin}  # Правый нижний
    ]
    
    # Добавляем скриншоты на слайд
    for screenshot, position in zip(screenshots, positions):
        if os.path.exists(screenshot):
            slide.shapes.add_picture(
                screenshot,
                left=position['left'],
                top=position['top'],
                width=width,
                height=height
            )
    
    prs.save('weather_report.pptx')
    print("Presentation created successfully")

async def main():
    try:
        if not os.path.exists('screenshots'):
            os.makedirs('screenshots')
            
        playwright = await async_playwright().start()
        browser = await playwright.chromium.launch(headless=False)
        context = await browser.new_context(
            viewport={'width': 1920, 'height': 1080},
            user_agent='Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        )
        page = await context.new_page()
        
        urls = [
            'https://www.gismeteo.ru/weather-moscow-4368/',
            'https://www.gismeteo.ru/weather-moscow-4368/tomorrow/',
            'https://www.gismeteo.ru/weather-moscow-4368/3-days/',
            'https://www.gismeteo.ru/weather-moscow-4368/10-days/'
        ]
        
        screenshots = []
        for i, url in enumerate(urls):
            filename = f'screenshots/weather_{i+1}.png'
            print(f"Taking screenshot of {url}")
            success = await take_screenshot(page, url, filename)
            if success:
                screenshots.append(filename)
            await page.wait_for_timeout(100)  # Reduced delay between requests to 0.1 second
        
        if screenshots:
            await create_presentation(screenshots)
            print("Presentation created successfully!")
        else:
            print("No screenshots were captured successfully.")
        
        await browser.close()
        await playwright.stop()
    except Exception as e:
        print(f"Произошла ошибка: {str(e)}")

if __name__ == "__main__":
    import asyncio
    asyncio.run(main()) 